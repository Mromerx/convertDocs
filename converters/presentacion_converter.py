import os
import shutil
from pathlib import Path
from .libreoffice_engine import convert_with_libreoffice

SUPPORTED_FORMATS = ['pptx', 'odp', 'pdf', 'txt']

def is_supported(ext: str):
    return ext.lower() in SUPPORTED_FORMATS

def convert(input_path: str, output_format: str, output_dir: str):
    """
    Converts a presentation to the specified format.
    """
    input_ext = Path(input_path).suffix.lower()[1:]
    output_format = output_format.lower()
    
    if input_ext == output_format:
        raise ValueError("Input and output format are the same.")
        
    if not is_supported(input_ext) or not is_supported(output_format):
        raise ValueError(f"Format not supported in presentation category. (Supported: {', '.join(SUPPORTED_FORMATS)})")
    
    output_file_name = f"{Path(input_path).stem}.{output_format}"
    output_path = os.path.join(output_dir, output_file_name)
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    
    # Text Extraction: PPTX -> TXT
    if input_ext == 'pptx' and output_format == 'txt':
        from pptx import Presentation
        prs = Presentation(input_path)
        with open(output_path, 'w', encoding='utf-8') as f:
            for i, slide in enumerate(prs.slides):
                f.write(f"--- Slide {i + 1} ---\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        f.write(shape.text + "\n")
                
                # Check for notes
                if slide.has_notes_slide:
                    text_frame = slide.notes_slide.notes_text_frame
                    if text_frame and text_frame.text.strip():
                        f.write("Speaker notes:\n")
                        f.write(text_frame.text + "\n")
                f.write("\n")
        return output_path
        
    # Text Generation: TXT -> PPTX
    elif input_ext == 'txt' and output_format in ('pptx', 'odp'):
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        # Using a blank layout
        blank_slide_layout = prs.slide_layouts[6]
        
        with open(input_path, 'r', encoding='utf-8') as f:
            content = f.read()
            # Split by double newline to create a slide for each block
            blocks = [b.strip() for b in content.split('\n\n') if b.strip()]
            
            for block in blocks:
                slide = prs.slides.add_slide(blank_slide_layout)
                left = top = width = height = Inches(1)
                # Just add a text box
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = block
                
        if output_format == 'pptx':
            prs.save(output_path)
            return output_path
        else:
            tmp_dir = os.path.join(output_dir, f"_tmp_pres_{Path(input_path).stem}")
            Path(tmp_dir).mkdir(parents=True, exist_ok=True)
            tmp_pptx = os.path.join(tmp_dir, f"{Path(input_path).stem}.pptx")
            prs.save(tmp_pptx)
            convert_with_libreoffice(tmp_pptx, 'odp', tmp_dir)
            tmp_odp = os.path.join(tmp_dir, f"{Path(input_path).stem}.odp")
            shutil.move(tmp_odp, output_path)
            shutil.rmtree(tmp_dir, ignore_errors=True)
            return output_path
        
    elif input_ext in ['pptx', 'odp', 'txt', 'pdf'] and output_format in ['pptx', 'odp', 'pdf', 'txt']:
        return convert_with_libreoffice(input_path, output_format, output_dir)
        
    else:
        if input_ext == 'pdf' and output_format == 'txt':
            raise ValueError("Use the generic PDF <-> TXT conversion")
        raise ValueError(f"Conversion {input_ext} -> {output_format} not implemented here.")
